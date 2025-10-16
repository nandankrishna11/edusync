"""
Text extraction utilities for PDF and PPTX files
"""

import os
from typing import Dict, List, Tuple
import PyPDF2
from pptx import Presentation
import logging

# Setup logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

def extract_text_from_pdf(file_path: str) -> Dict[str, any]:
    """
    Extract text from PDF file with page information
    
    Args:
        file_path (str): Path to the PDF file
        
    Returns:
        Dict containing:
        - text (str): Full extracted text
        - pages (List[Dict]): List of page data with text and page number
        - metadata (Dict): File metadata
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"PDF file not found: {file_path}")
    
    if not file_path.lower().endswith('.pdf'):
        raise ValueError("File must be a PDF")
    
    try:
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            
            # Extract metadata
            metadata = {
                'num_pages': len(pdf_reader.pages),
                'file_name': os.path.basename(file_path),
                'file_size': os.path.getsize(file_path)
            }
            
            # Try to get PDF metadata if available
            if pdf_reader.metadata:
                metadata.update({
                    'title': pdf_reader.metadata.get('/Title', ''),
                    'author': pdf_reader.metadata.get('/Author', ''),
                    'subject': pdf_reader.metadata.get('/Subject', ''),
                    'creator': pdf_reader.metadata.get('/Creator', '')
                })
            
            # Extract text from each page
            pages = []
            full_text = ""
            
            for page_num, page in enumerate(pdf_reader.pages, 1):
                try:
                    page_text = page.extract_text()
                    
                    # Clean up the text
                    page_text = page_text.strip()
                    
                    if page_text:  # Only add non-empty pages
                        pages.append({
                            'page_number': page_num,
                            'text': page_text,
                            'char_count': len(page_text)
                        })
                        full_text += f"\\n\\n--- Page {page_num} ---\\n\\n{page_text}"
                    
                except Exception as e:
                    logger.warning(f"Error extracting text from page {page_num}: {e}")
                    pages.append({
                        'page_number': page_num,
                        'text': '',
                        'char_count': 0,
                        'error': str(e)
                    })
            
            result = {
                'text': full_text.strip(),
                'pages': pages,
                'metadata': metadata,
                'extraction_success': True
            }
            
            logger.info(f"Successfully extracted text from PDF: {metadata['num_pages']} pages, {len(full_text)} characters")
            return result
            
    except Exception as e:
        logger.error(f"Error extracting text from PDF {file_path}: {e}")
        return {
            'text': '',
            'pages': [],
            'metadata': {'file_name': os.path.basename(file_path), 'error': str(e)},
            'extraction_success': False,
            'error': str(e)
        }

def extract_text_from_pptx(file_path: str) -> Dict[str, any]:
    """
    Extract text from PPTX file with slide information
    
    Args:
        file_path (str): Path to the PPTX file
        
    Returns:
        Dict containing:
        - text (str): Full extracted text
        - slides (List[Dict]): List of slide data with text and slide number
        - metadata (Dict): File metadata
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"PPTX file not found: {file_path}")
    
    if not file_path.lower().endswith(('.pptx', '.ppt')):
        raise ValueError("File must be a PowerPoint file (.pptx or .ppt)")
    
    try:
        presentation = Presentation(file_path)
        
        # Extract metadata
        metadata = {
            'num_slides': len(presentation.slides),
            'file_name': os.path.basename(file_path),
            'file_size': os.path.getsize(file_path)
        }
        
        # Try to get presentation metadata if available
        if hasattr(presentation.core_properties, 'title') and presentation.core_properties.title:
            metadata['title'] = presentation.core_properties.title
        if hasattr(presentation.core_properties, 'author') and presentation.core_properties.author:
            metadata['author'] = presentation.core_properties.author
        if hasattr(presentation.core_properties, 'subject') and presentation.core_properties.subject:
            metadata['subject'] = presentation.core_properties.subject
        
        # Extract text from each slide
        slides = []
        full_text = ""
        
        for slide_num, slide in enumerate(presentation.slides, 1):
            try:
                slide_text = ""
                
                # Extract text from all shapes in the slide
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text += shape.text + "\\n"
                
                # Clean up the text
                slide_text = slide_text.strip()
                
                if slide_text:  # Only add non-empty slides
                    slides.append({
                        'slide_number': slide_num,
                        'text': slide_text,
                        'char_count': len(slide_text)
                    })
                    full_text += f"\\n\\n--- Slide {slide_num} ---\\n\\n{slide_text}"
                else:
                    # Add empty slide info
                    slides.append({
                        'slide_number': slide_num,
                        'text': '',
                        'char_count': 0,
                        'note': 'No text content found'
                    })
                
            except Exception as e:
                logger.warning(f"Error extracting text from slide {slide_num}: {e}")
                slides.append({
                    'slide_number': slide_num,
                    'text': '',
                    'char_count': 0,
                    'error': str(e)
                })
        
        result = {
            'text': full_text.strip(),
            'slides': slides,
            'metadata': metadata,
            'extraction_success': True
        }
        
        logger.info(f"Successfully extracted text from PPTX: {metadata['num_slides']} slides, {len(full_text)} characters")
        return result
        
    except Exception as e:
        logger.error(f"Error extracting text from PPTX {file_path}: {e}")
        return {
            'text': '',
            'slides': [],
            'metadata': {'file_name': os.path.basename(file_path), 'error': str(e)},
            'extraction_success': False,
            'error': str(e)
        }

def get_supported_file_types() -> List[str]:
    """
    Get list of supported file types for text extraction
    
    Returns:
        List[str]: Supported file extensions
    """
    return ['.pdf', '.pptx', '.ppt']

def is_supported_file(file_path: str) -> bool:
    """
    Check if file type is supported for text extraction
    
    Args:
        file_path (str): Path to the file
        
    Returns:
        bool: True if file type is supported
    """
    return any(file_path.lower().endswith(ext) for ext in get_supported_file_types())